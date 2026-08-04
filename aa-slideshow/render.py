"""Render a .pptx to HTML at true geometry (read back from the file itself),
so it can be screenshotted with Chromium. LibreOffice is broken in this sandbox.
"""
import sys, html
from pptx import Presentation
from pptx.util import Emu

PX = 100.0  # px per inch

def inches(v):
    return Emu(v).inches if v is not None else 0.0

def color_of(cf):
    try:
        if cf and cf.type is not None and cf.rgb is not None:
            return "#%s" % str(cf.rgb)
    except Exception:
        pass
    return None

def shape_fill(sp):
    try:
        f = sp.fill
        if f.type is not None and f.type == 1:  # solid
            return "#%s" % str(f.fore_color.rgb)
    except Exception:
        pass
    return None

def main(path, out):
    prs = Presentation(path)
    W, H = prs.slide_width.inches * PX, prs.slide_height.inches * PX
    parts = ["""<!doctype html><meta charset=utf-8><style>
    body{margin:0;background:#888}
    .slide{position:relative;width:%dpx;height:%dpx;background:#fff;overflow:hidden}
    .b{position:absolute;box-sizing:border-box}
    .t{position:absolute;display:flex;box-sizing:border-box;line-height:1.2}
    .t>span{width:100%%}
    </style>""" % (W, H)]

    for si, slide in enumerate(prs.slides):
        parts.append('<div class="slide">')
        for sp in slide.shapes:
            x, y = inches(sp.left) * PX, inches(sp.top) * PX
            w, h = inches(sp.width) * PX, inches(sp.height) * PX
            fill = shape_fill(sp)
            if fill and not (sp.has_text_frame and sp.text_frame.text.strip()):
                parts.append('<div class="b" style="left:%.1fpx;top:%.1fpx;width:%.1fpx;height:%.1fpx;'
                             'background:%s;border-radius:8px"></div>' % (x, y, w, h, fill))
            elif fill:
                parts.append('<div class="b" style="left:%.1fpx;top:%.1fpx;width:%.1fpx;height:%.1fpx;'
                             'background:%s;border-radius:8px"></div>' % (x, y, w, h, fill))
            if not sp.has_text_frame:
                continue
            tf = sp.text_frame
            if not tf.text.strip():
                continue
            va = {None: "flex-start", 1: "center", 2: "flex-end"}.get(
                getattr(tf.vertical_anchor, "value", None) if tf.vertical_anchor is not None else None,
                "center" if tf.vertical_anchor is not None else "flex-start")
            try:
                va = "center" if str(tf.vertical_anchor) .startswith("MIDDLE") else va
            except Exception:
                pass
            inner = []
            for p in tf.paragraphs:
                al = str(p.alignment or "")
                align = "right" if "RIGHT" in al else ("center" if "CENTER" in al else "left")
                runs = []
                for r in p.runs:
                    st = []
                    fnt = r.font
                    if fnt.size: st.append("font-size:%.1fpx" % (fnt.size.pt * PX / 72.0))
                    if fnt.bold: st.append("font-weight:700")
                    nm = (fnt.name or "")
                    st.append("font-family:%s" % ("Georgia,serif" if "Cambria" in nm else "Helvetica,Arial,sans-serif"))
                    c = color_of(fnt.color)
                    if c: st.append("color:%s" % c)
                    strike = r._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                    if strike is not None and strike.get('strike') == 'sngStrike':
                        st.append("text-decoration:line-through")
                    runs.append('<span style="%s">%s</span>' % (";".join(st), html.escape(r.text)))
                inner.append('<div style="text-align:%s">%s</div>' % (align, "".join(runs) or "&nbsp;"))
            parts.append('<div class="t" style="left:%.1fpx;top:%.1fpx;width:%.1fpx;height:%.1fpx;'
                         'align-items:%s"><span>%s</span></div>'
                         % (x, y, w, h, va, "".join(inner)))
        parts.append("</div>")
    open(out, "w").write("".join(parts))
    print("wrote", out, "%dx%d" % (W, H))

if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2])
