"""Rebuild pages/ from the committed section decks.

The rendered page images are gitignored, so a fresh container has the decks but
not the JPEGs that went into them. Every image is still inside the .pptx files:
walk the slides in order, take each spread's pictures left to right, and write
them back out under the sequence numbers recorded in pages.json.
"""
import json, os
from pptx import Presentation

SECTIONS = [
 ("5215354a-08182026_ALCOHOLIC_S_ANONYMOUS.pdf", "big-book-1-front-matter.pptx"),
 ("43630d84-08182026_THERE_IS__SOLUTION_D_Njj__t7_i5uri_Cr__ft_VG_Lt_6TLi.pdf", "big-book-2-there-is-a-solution.pptx"),
 ("231419ad-08182026___INTO_ACTION_73_invariably_they_got_drunk._Having_per.pdf", "big-book-3-into-action.pptx"),
 ("6c79fd5d-08182026_TO_EMPLOYERS_141_normal_will_do_incredible_things._Aft.pdf", "big-book-4-to-employers.pptx"),
]

pages = json.load(open("pages.json"))
os.makedirs("pages", exist_ok=True)
written = 0
for src, deck in SECTIONS:
    want = [p for p in pages if p["source"] == src]
    prs = Presentation(deck)
    got = []
    for slide in prs.slides:
        pics = sorted([sh for sh in slide.shapes if sh.shape_type == 13], key=lambda sh: sh.left)
        for sh in pics:
            got.append(sh.image.blob)
    if len(got) != len(want):
        raise SystemExit("%s: %d images but %d pages expected" % (deck, len(got), len(want)))
    for p, blob in zip(want, got):
        open(os.path.join("pages", p["file"]), "wb").write(blob)
        written += 1
    print("  %-40s %3d pages" % (deck, len(want)))
print("recovered %d page images" % written)
